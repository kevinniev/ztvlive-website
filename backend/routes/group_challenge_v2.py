"""
ZTVLIVE Private Group Challenge System - Enhanced Edition

Features:
- Create private groups with passcode
- Invite via email, passcode, or social share
- Custom question import (PPT, Word, Excel - auto-detect)
- Jitsi video chat with presenter mode
- Host/Co-host game controls (next, play, pause)
- Dual leaderboard: Group Top 4 + Global Top 4
"""

from fastapi import APIRouter, HTTPException, UploadFile, File, Form, WebSocket, WebSocketDisconnect
from pydantic import BaseModel, EmailStr
from typing import Dict, List, Optional, Any
from datetime import datetime, timezone
import asyncio
import random
import string
import uuid
import os
import io
import re
from motor.motor_asyncio import AsyncIOMotorClient

# Document parsing libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

router = APIRouter(prefix="/api/game/groups", tags=["Group Challenge"])

# MongoDB connection
MONGO_URL = os.environ.get("MONGO_URL", "mongodb://localhost:27017")
DB_NAME = os.environ.get("DB_NAME", "ztvlive")
client = AsyncIOMotorClient(MONGO_URL)
db = client[DB_NAME]

# ============== MODELS ==============

class CreateGroupRequest(BaseModel):
    name: str
    host_name: str
    host_email: Optional[str] = None
    is_permanent: bool = False
    max_members: int = 50
    enable_video: bool = True
    enable_passcode: bool = True

class JoinGroupRequest(BaseModel):
    player_name: str
    player_email: Optional[str] = None
    passcode: Optional[str] = None

class InviteMembersRequest(BaseModel):
    emails: List[EmailStr]
    custom_message: Optional[str] = None

class AddCoHostRequest(BaseModel):
    member_id: str

class GameControlRequest(BaseModel):
    action: str  # "next", "pause", "resume", "end", "show_results"
    question_index: Optional[int] = None

class SpotlightRequest(BaseModel):
    member_id: str  # Member to spotlight (or "host" for default)

class SubmitAnswerRequest(BaseModel):
    player_id: str
    answer: str
    question_id: str

# ============== PASSCODE GENERATOR ==============

def generate_passcode(length: int = 6) -> str:
    """Generate a simple numeric passcode"""
    return ''.join(random.choices(string.digits, k=length))

def generate_group_id() -> str:
    """Generate a short memorable group ID"""
    return ''.join(random.choices(string.ascii_uppercase + string.digits, k=8))

# ============== QUESTION PARSER ==============

class QuestionParser:
    """Auto-detect and parse questions from PPT, Word, or Excel files"""
    
    @staticmethod
    def parse_pptx(file_bytes: bytes) -> List[dict]:
        """Parse questions from PowerPoint file"""
        if not PPTX_AVAILABLE:
            raise HTTPException(status_code=500, detail="PPTX parsing not available")
        
        questions = []
        prs = Presentation(io.BytesIO(file_bytes))
        
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            
            if texts:
                # First text is usually the question, rest are options
                question_text = texts[0] if texts else f"Question {slide_num}"
                options = texts[1:5] if len(texts) > 1 else []
                
                # Try to detect if it's a question (ends with ?)
                if '?' in question_text or len(options) >= 2:
                    questions.append({
                        "id": str(uuid.uuid4())[:8],
                        "question": question_text,
                        "options": options[:4],  # Max 4 options
                        "type": "multiple_choice" if options else "open_ended",
                        "source_slide": slide_num
                    })
        
        return questions
    
    @staticmethod
    def parse_docx(file_bytes: bytes) -> List[dict]:
        """Parse questions from Word document"""
        if not DOCX_AVAILABLE:
            raise HTTPException(status_code=500, detail="DOCX parsing not available")
        
        questions = []
        doc = Document(io.BytesIO(file_bytes))
        
        current_question = None
        current_options = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Detect question patterns
            # Pattern 1: "Q1. What is...?" or "1. What is...?"
            q_match = re.match(r'^[Qq]?\d+[\.\)]\s*(.+)', text)
            # Pattern 2: Line ending with ?
            is_question = text.endswith('?') or q_match
            # Pattern 3: Options like "A) ..." or "a. ..."
            opt_match = re.match(r'^[A-Da-d][\.\)]\s*(.+)', text)
            
            if is_question and not opt_match:
                # Save previous question if exists
                if current_question:
                    questions.append({
                        "id": str(uuid.uuid4())[:8],
                        "question": current_question,
                        "options": current_options[:4],
                        "type": "multiple_choice" if current_options else "open_ended"
                    })
                current_question = q_match.group(1) if q_match else text
                current_options = []
            elif opt_match:
                current_options.append(opt_match.group(1))
            elif current_question and len(current_options) < 4:
                # Might be an option without letter prefix
                current_options.append(text)
        
        # Don't forget last question
        if current_question:
            questions.append({
                "id": str(uuid.uuid4())[:8],
                "question": current_question,
                "options": current_options[:4],
                "type": "multiple_choice" if current_options else "open_ended"
            })
        
        return questions
    
    @staticmethod
    def parse_xlsx(file_bytes: bytes) -> List[dict]:
        """Parse questions from Excel file"""
        if not XLSX_AVAILABLE:
            raise HTTPException(status_code=500, detail="XLSX parsing not available")
        
        questions = []
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
        ws = wb.active
        
        # Try to detect header row
        headers = []
        for cell in ws[1]:
            if cell.value:
                headers.append(str(cell.value).lower())
        
        # Common header patterns
        q_col = None
        opt_cols = []
        
        for i, h in enumerate(headers):
            if 'question' in h or 'q ' in h or h == 'q':
                q_col = i
            elif 'option' in h or 'answer' in h or 'choice' in h or h in ['a', 'b', 'c', 'd']:
                opt_cols.append(i)
        
        # If no headers detected, assume: Col A = Question, B-E = Options
        if q_col is None:
            q_col = 0
            opt_cols = [1, 2, 3, 4] if ws.max_column > 1 else []
        
        # Parse rows (skip header)
        for row_num, row in enumerate(ws.iter_rows(min_row=2, values_only=True), 2):
            if not row or not row[q_col]:
                continue
            
            question_text = str(row[q_col]).strip()
            options = []
            
            for col_idx in opt_cols:
                if col_idx < len(row) and row[col_idx]:
                    options.append(str(row[col_idx]).strip())
            
            if question_text:
                questions.append({
                    "id": str(uuid.uuid4())[:8],
                    "question": question_text,
                    "options": options[:4],
                    "type": "multiple_choice" if options else "open_ended",
                    "source_row": row_num
                })
        
        return questions
    
    @staticmethod
    def auto_parse(filename: str, file_bytes: bytes) -> List[dict]:
        """Auto-detect file type and parse questions"""
        ext = filename.lower().split('.')[-1]
        
        if ext in ['pptx', 'ppt']:
            return QuestionParser.parse_pptx(file_bytes)
        elif ext in ['docx', 'doc']:
            return QuestionParser.parse_docx(file_bytes)
        elif ext in ['xlsx', 'xls', 'csv']:
            return QuestionParser.parse_xlsx(file_bytes)
        else:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file type: {ext}. Supported: pptx, docx, xlsx"
            )


# ============== IN-MEMORY GROUP STATE ==============

class GroupState:
    """Real-time state for an active group session"""
    def __init__(self, group_id: str):
        self.group_id = group_id
        self.connected_clients: Dict[str, WebSocket] = {}
        self.members: Dict[str, dict] = {}
        self.co_hosts: set = set()  # player_ids of co-hosts
        self.scores: Dict[str, int] = {}
        self.current_answers: Dict[str, dict] = {}
        
        # Game state
        self.game_status: str = "lobby"  # lobby, playing, paused, ended
        self.current_question_index: int = -1
        self.custom_questions: List[dict] = []
        self.using_custom_questions: bool = False
        
        # Video state
        self.video_enabled: bool = True
        self.spotlight_member: Optional[str] = None  # player_id being spotlighted
        self.jitsi_room_name: Optional[str] = None
        
    def add_member(self, player_id: str, member_data: dict):
        self.members[player_id] = member_data
        self.scores[player_id] = member_data.get("score", 0)
        
    def remove_member(self, player_id: str):
        self.members.pop(player_id, None)
        self.scores.pop(player_id, None)
        self.current_answers.pop(player_id, None)
        self.co_hosts.discard(player_id)
        if player_id in self.connected_clients:
            del self.connected_clients[player_id]
            
    def add_co_host(self, player_id: str):
        if player_id in self.members:
            self.co_hosts.add(player_id)
            self.members[player_id]["is_co_host"] = True
            
    def is_host_or_cohost(self, player_id: str) -> bool:
        member = self.members.get(player_id, {})
        return member.get("is_host", False) or player_id in self.co_hosts
        
    def get_current_question(self) -> Optional[dict]:
        if self.using_custom_questions and 0 <= self.current_question_index < len(self.custom_questions):
            return self.custom_questions[self.current_question_index]
        return None
        
    def next_question(self) -> Optional[dict]:
        if self.using_custom_questions:
            self.current_question_index += 1
            self.current_answers.clear()  # Clear answers for new question
            return self.get_current_question()
        return None
        
    def record_answer(self, player_id: str, answer: str, question_id: str):
        self.current_answers[player_id] = {
            "answer": answer,
            "question_id": question_id,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "player_name": self.members.get(player_id, {}).get("name", "Anonymous")
        }
        
    def get_group_leaderboard(self) -> List[dict]:
        """Get top 4 players in the group"""
        sorted_scores = sorted(
            [(pid, self.scores.get(pid, 0)) for pid in self.members.keys()],
            key=lambda x: x[1],
            reverse=True
        )[:4]
        
        return [
            {
                "rank": i + 1,
                "player_id": pid,
                "player_name": self.members.get(pid, {}).get("name", "Anonymous"),
                "score": score,
                "is_host": self.members.get(pid, {}).get("is_host", False),
                "is_co_host": pid in self.co_hosts
            }
            for i, (pid, score) in enumerate(sorted_scores)
        ]
        
    def get_live_answers(self) -> List[dict]:
        return [
            {
                "player_id": pid,
                "player_name": data.get("player_name", "Anonymous"),
                "answer": data.get("answer"),
                "timestamp": data.get("timestamp")
            }
            for pid, data in self.current_answers.items()
        ]


# Active group sessions
active_groups: Dict[str, GroupState] = {}

def get_or_create_group_state(group_id: str) -> GroupState:
    if group_id not in active_groups:
        active_groups[group_id] = GroupState(group_id)
    return active_groups[group_id]


# ============== API ENDPOINTS ==============

@router.post("/create")
async def create_group(request: CreateGroupRequest):
    """Create a new private group with passcode"""
    group_id = generate_group_id()
    passcode = generate_passcode() if request.enable_passcode else None
    invite_code = str(uuid.uuid4())[:12]
    host_id = str(uuid.uuid4())
    jitsi_room = f"ztvlive-{group_id.lower()}" if request.enable_video else None
    
    group_data = {
        "group_id": group_id,
        "passcode": passcode,
        "invite_code": invite_code,
        "name": request.name,
        "host_id": host_id,
        "host_name": request.host_name,
        "host_email": request.host_email,
        "is_permanent": request.is_permanent,
        "max_members": request.max_members,
        "enable_video": request.enable_video,
        "jitsi_room": jitsi_room,
        "co_hosts": [],
        "members": [{
            "player_id": host_id,
            "name": request.host_name,
            "email": request.host_email,
            "is_host": True,
            "is_co_host": False,
            "joined_at": datetime.now(timezone.utc),
            "score": 0
        }],
        "custom_questions": [],
        "created_at": datetime.now(timezone.utc),
        "status": "active",
        "game_state": "lobby"
    }
    
    await db.game_groups.insert_one(group_data)
    
    # Initialize in-memory state
    group_state = get_or_create_group_state(group_id)
    group_state.add_member(host_id, {
        "name": request.host_name,
        "email": request.host_email,
        "is_host": True,
        "score": 0
    })
    group_state.video_enabled = request.enable_video
    group_state.jitsi_room_name = jitsi_room
    
    base_url = os.environ.get('BASE_URL', 'https://www.ztvlivestream.com')
    
    return {
        "success": True,
        "group_id": group_id,
        "passcode": passcode,
        "invite_code": invite_code,
        "invite_link": f"{base_url}/play?group={group_id}",
        "host_id": host_id,
        "jitsi_room": jitsi_room,
        "social_share": {
            "text": f"Join my ZTVLIVE UNUSUAL FUN game! Code: {passcode}",
            "url": f"{base_url}/play?group={group_id}",
            "hashtags": "ZTVLIVE,UnusualFun,GameNight,Trivia"
        }
    }


@router.get("/{group_id}")
async def get_group(group_id: str):
    """Get group details"""
    group = await db.game_groups.find_one({"group_id": group_id}, {"_id": 0})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    group_state = active_groups.get(group_id)
    
    return {
        **group,
        "live_member_count": len(group_state.members) if group_state else len(group.get("members", [])),
        "leaderboard": group_state.get_group_leaderboard() if group_state else [],
        "game_status": group_state.game_status if group_state else "lobby",
        "current_question": group_state.get_current_question() if group_state else None,
        "using_custom_questions": group_state.using_custom_questions if group_state else False
    }


@router.post("/{group_id}/join")
async def join_group(group_id: str, request: JoinGroupRequest):
    """Join a group via invite link or passcode"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    # Verify passcode if provided
    if request.passcode:
        if group.get("passcode") != request.passcode:
            raise HTTPException(status_code=403, detail="Invalid passcode")
    
    if len(group.get("members", [])) >= group.get("max_members", 50):
        raise HTTPException(status_code=400, detail="Group is full")
    
    # Check if already in group
    existing = next(
        (m for m in group.get("members", []) 
         if m.get("email") == request.player_email and request.player_email),
        None
    )
    if existing:
        return {
            "success": True,
            "player_id": existing["player_id"],
            "message": "Welcome back!",
            "group_name": group["name"],
            "jitsi_room": group.get("jitsi_room")
        }
    
    player_id = str(uuid.uuid4())
    
    member_data = {
        "player_id": player_id,
        "name": request.player_name,
        "email": request.player_email,
        "is_host": False,
        "is_co_host": False,
        "is_guest": request.player_email is None,
        "joined_at": datetime.now(timezone.utc),
        "score": 0
    }
    
    await db.game_groups.update_one(
        {"group_id": group_id},
        {"$push": {"members": member_data}}
    )
    
    group_state = get_or_create_group_state(group_id)
    group_state.add_member(player_id, {
        "name": request.player_name,
        "email": request.player_email,
        "is_host": False,
        "score": 0
    })
    
    await broadcast_to_group(group_id, {
        "type": "member_joined",
        "player_id": player_id,
        "player_name": request.player_name,
        "member_count": len(group_state.members)
    })
    
    return {
        "success": True,
        "player_id": player_id,
        "group_name": group["name"],
        "jitsi_room": group.get("jitsi_room"),
        "game_status": group_state.game_status,
        "using_custom_questions": group_state.using_custom_questions
    }


@router.post("/{group_id}/join-by-passcode")
async def join_by_passcode(group_id: str, passcode: str, player_name: str):
    """Quick join via passcode only"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    if group.get("passcode") != passcode:
        raise HTTPException(status_code=403, detail="Invalid passcode")
    
    # Same as regular join
    request = JoinGroupRequest(player_name=player_name, passcode=passcode)
    return await join_group(group_id, request)


@router.post("/{group_id}/upload-questions")
async def upload_questions(
    group_id: str,
    host_id: str = Form(...),
    file: UploadFile = File(...)
):
    """Upload custom questions from PPT, Word, or Excel"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    # Verify host or co-host
    group_state = active_groups.get(group_id)
    if group_state and not group_state.is_host_or_cohost(host_id):
        if group.get("host_id") != host_id:
            raise HTTPException(status_code=403, detail="Only host/co-host can upload questions")
    
    # Parse file
    file_bytes = await file.read()
    questions = QuestionParser.auto_parse(file.filename, file_bytes)
    
    if not questions:
        raise HTTPException(status_code=400, detail="No questions found in file")
    
    # Store questions
    await db.game_groups.update_one(
        {"group_id": group_id},
        {"$set": {"custom_questions": questions}}
    )
    
    # Update live state
    if group_state:
        group_state.custom_questions = questions
        group_state.using_custom_questions = True
        group_state.current_question_index = -1
        
        await broadcast_to_group(group_id, {
            "type": "questions_uploaded",
            "question_count": len(questions),
            "first_question": questions[0] if questions else None
        })
    
    return {
        "success": True,
        "questions_parsed": len(questions),
        "questions": questions,
        "message": f"Successfully imported {len(questions)} questions from {file.filename}"
    }


@router.post("/{group_id}/add-cohost")
async def add_co_host(group_id: str, request: AddCoHostRequest, host_id: str):
    """Add a co-host (host only)"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    if group.get("host_id") != host_id:
        raise HTTPException(status_code=403, detail="Only the host can add co-hosts")
    
    # Update database
    await db.game_groups.update_one(
        {"group_id": group_id},
        {
            "$addToSet": {"co_hosts": request.member_id},
            "$set": {"members.$[m].is_co_host": True}
        },
        array_filters=[{"m.player_id": request.member_id}]
    )
    
    # Update live state
    group_state = active_groups.get(group_id)
    if group_state:
        group_state.add_co_host(request.member_id)
        
        await broadcast_to_group(group_id, {
            "type": "co_host_added",
            "player_id": request.member_id,
            "player_name": group_state.members.get(request.member_id, {}).get("name")
        })
    
    return {"success": True, "message": "Co-host added"}


@router.post("/{group_id}/game-control")
async def game_control(group_id: str, request: GameControlRequest, controller_id: str):
    """Host/Co-host game controls: next, pause, resume, end"""
    group_state = active_groups.get(group_id)
    
    if not group_state:
        raise HTTPException(status_code=404, detail="Group session not found")
    
    if not group_state.is_host_or_cohost(controller_id):
        raise HTTPException(status_code=403, detail="Only host/co-host can control the game")
    
    action = request.action.lower()
    response_data = {"success": True, "action": action}
    
    if action == "start":
        group_state.game_status = "playing"
        if group_state.using_custom_questions:
            group_state.current_question_index = 0
        response_data["game_status"] = "playing"
        response_data["current_question"] = group_state.get_current_question()
        
    elif action == "next":
        question = group_state.next_question()
        response_data["current_question"] = question
        response_data["question_index"] = group_state.current_question_index
        if question is None and group_state.using_custom_questions:
            response_data["message"] = "No more questions"
            group_state.game_status = "ended"
            
    elif action == "pause":
        group_state.game_status = "paused"
        response_data["game_status"] = "paused"
        
    elif action == "resume":
        group_state.game_status = "playing"
        response_data["game_status"] = "playing"
        
    elif action == "end":
        group_state.game_status = "ended"
        response_data["game_status"] = "ended"
        response_data["final_leaderboard"] = group_state.get_group_leaderboard()
        
    elif action == "show_results":
        response_data["live_answers"] = group_state.get_live_answers()
        response_data["leaderboard"] = group_state.get_group_leaderboard()
    
    # Broadcast to all members
    await broadcast_to_group(group_id, {
        "type": "game_control",
        **response_data
    })
    
    return response_data


@router.post("/{group_id}/spotlight")
async def set_spotlight(group_id: str, request: SpotlightRequest, controller_id: str):
    """Set which member's video is spotlighted (presenter mode)"""
    group_state = active_groups.get(group_id)
    
    if not group_state:
        raise HTTPException(status_code=404, detail="Group session not found")
    
    if not group_state.is_host_or_cohost(controller_id):
        raise HTTPException(status_code=403, detail="Only host/co-host can set spotlight")
    
    group_state.spotlight_member = request.member_id
    
    await broadcast_to_group(group_id, {
        "type": "spotlight_changed",
        "spotlight_member": request.member_id,
        "spotlight_name": group_state.members.get(request.member_id, {}).get("name", "Host")
    })
    
    return {
        "success": True,
        "spotlight_member": request.member_id
    }


@router.post("/{group_id}/answer")
async def submit_answer(group_id: str, request: SubmitAnswerRequest):
    """Submit an answer"""
    group_state = active_groups.get(group_id)
    
    if not group_state:
        raise HTTPException(status_code=404, detail="Group session not found")
    
    group_state.record_answer(request.player_id, request.answer, request.question_id)
    
    await broadcast_to_group(group_id, {
        "type": "member_answered",
        "player_id": request.player_id,
        "player_name": group_state.members.get(request.player_id, {}).get("name"),
        "answer": request.answer,
        "live_answers": group_state.get_live_answers()
    })
    
    return {"success": True}


@router.post("/{group_id}/update-score")
async def update_score(group_id: str, player_id: str, points: int):
    """Update player score"""
    group_state = active_groups.get(group_id)
    
    if not group_state or player_id not in group_state.members:
        return {"success": False}
    
    group_state.scores[player_id] = group_state.scores.get(player_id, 0) + points
    
    await db.game_groups.update_one(
        {"group_id": group_id, "members.player_id": player_id},
        {"$inc": {"members.$.score": points}}
    )
    
    await broadcast_to_group(group_id, {
        "type": "score_update",
        "player_id": player_id,
        "new_score": group_state.scores[player_id],
        "leaderboard": group_state.get_group_leaderboard()
    })
    
    return {"success": True, "new_score": group_state.scores[player_id]}


@router.get("/{group_id}/leaderboard")
async def get_leaderboard(group_id: str):
    """Get group leaderboard"""
    group_state = active_groups.get(group_id)
    
    if not group_state:
        group = await db.game_groups.find_one({"group_id": group_id})
        if not group:
            raise HTTPException(status_code=404, detail="Group not found")
        
        members = sorted(group.get("members", []), key=lambda m: m.get("score", 0), reverse=True)[:4]
        return {
            "group_id": group_id,
            "leaderboard": [
                {"rank": i+1, "player_name": m["name"], "score": m.get("score", 0)}
                for i, m in enumerate(members)
            ]
        }
    
    return {
        "group_id": group_id,
        "leaderboard": group_state.get_group_leaderboard(),
        "live_answers": group_state.get_live_answers()
    }


@router.post("/{group_id}/invite")
async def invite_members(group_id: str, request: InviteMembersRequest, host_id: str):
    """Send email invites"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    group_state = active_groups.get(group_id)
    if group_state and not group_state.is_host_or_cohost(host_id):
        if group.get("host_id") != host_id:
            raise HTTPException(status_code=403, detail="Only host/co-host can send invites")
    
    base_url = os.environ.get('BASE_URL', 'https://www.ztvlivestream.com')
    invite_link = f"{base_url}/play?group={group_id}"
    passcode = group.get("passcode", "")
    
    from services.email_service import send_email
    
    sent_count = 0
    for email in request.emails:
        try:
            html = get_invite_email_html(
                group["name"], group["host_name"], invite_link, passcode, request.custom_message
            )
            await send_email(email, f"🎮 Join {group['host_name']}'s ZTVLIVE Game!", html)
            sent_count += 1
        except Exception:
            pass
    
    return {"success": True, "sent_count": sent_count, "invite_link": invite_link}


@router.post("/{group_id}/kick")
async def kick_member(group_id: str, member_id: str, host_id: str):
    """Kick a member"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group:
        raise HTTPException(status_code=404, detail="Group not found")
    
    group_state = active_groups.get(group_id)
    if group_state and not group_state.is_host_or_cohost(host_id):
        if group.get("host_id") != host_id:
            raise HTTPException(status_code=403, detail="Only host/co-host can kick members")
    
    if member_id == group.get("host_id"):
        raise HTTPException(status_code=400, detail="Cannot kick the host")
    
    await db.game_groups.update_one(
        {"group_id": group_id},
        {"$pull": {"members": {"player_id": member_id}}}
    )
    
    if group_state:
        member_name = group_state.members.get(member_id, {}).get("name", "Unknown")
        
        if member_id in group_state.connected_clients:
            try:
                await group_state.connected_clients[member_id].send_json({
                    "type": "kicked",
                    "message": "You have been removed from the group"
                })
            except Exception:
                pass
        
        group_state.remove_member(member_id)
        
        await broadcast_to_group(group_id, {
            "type": "member_kicked",
            "player_id": member_id,
            "player_name": member_name
        })
    
    return {"success": True}


@router.delete("/{group_id}")
async def close_group(group_id: str, host_id: str):
    """Close group"""
    group = await db.game_groups.find_one({"group_id": group_id})
    
    if not group or group.get("host_id") != host_id:
        raise HTTPException(status_code=403, detail="Only host can close group")
    
    await db.game_groups.update_one(
        {"group_id": group_id},
        {"$set": {"status": "closed"}}
    )
    
    group_state = active_groups.get(group_id)
    if group_state:
        await broadcast_to_group(group_id, {"type": "group_closed"})
        for ws in list(group_state.connected_clients.values()):
            try:
                await ws.close()
            except Exception:
                pass
        del active_groups[group_id]
    
    return {"success": True}


# ============== WEBSOCKET ==============

@router.websocket("/{group_id}/ws")
async def group_websocket(websocket: WebSocket, group_id: str, player_id: str):
    """WebSocket for real-time updates"""
    await websocket.accept()
    
    group_state = get_or_create_group_state(group_id)
    group_state.connected_clients[player_id] = websocket
    
    try:
        await websocket.send_json({
            "type": "connected",
            "group_id": group_id,
            "game_status": group_state.game_status,
            "leaderboard": group_state.get_group_leaderboard(),
            "current_question": group_state.get_current_question(),
            "spotlight_member": group_state.spotlight_member,
            "jitsi_room": group_state.jitsi_room_name
        })
        
        while True:
            data = await websocket.receive_json()
            
            if data.get("type") == "answer":
                group_state.record_answer(player_id, data.get("answer"), data.get("question_id", ""))
                await broadcast_to_group(group_id, {
                    "type": "member_answered",
                    "player_id": player_id,
                    "player_name": group_state.members.get(player_id, {}).get("name"),
                    "answer": data.get("answer"),
                    "live_answers": group_state.get_live_answers()
                })
            elif data.get("type") == "ping":
                await websocket.send_json({"type": "pong"})
                
    except WebSocketDisconnect:
        group_state.connected_clients.pop(player_id, None)
        await broadcast_to_group(group_id, {
            "type": "member_disconnected",
            "player_id": player_id,
            "online_count": len(group_state.connected_clients)
        })


async def broadcast_to_group(group_id: str, message: dict):
    """Broadcast to all connected members"""
    group_state = active_groups.get(group_id)
    if not group_state:
        return
    
    disconnected = []
    for pid, ws in group_state.connected_clients.items():
        try:
            await ws.send_json(message)
        except Exception:
            disconnected.append(pid)
    
    for pid in disconnected:
        group_state.connected_clients.pop(pid, None)


# ============== EMAIL TEMPLATE ==============

def get_invite_email_html(group_name: str, host_name: str, invite_link: str, passcode: str, custom_message: str = None) -> str:
    custom_section = f'<p style="color:#a3a3a3;font-style:italic;">"{custom_message}"</p>' if custom_message else ""
    passcode_section = f'<div style="background:#262626;padding:15px;border-radius:8px;text-align:center;margin:20px 0;"><p style="color:#737373;margin:0 0 5px;">PASSCODE</p><h2 style="color:#10b981;margin:0;font-size:32px;letter-spacing:4px;">{passcode}</h2></div>' if passcode else ""
    
    return f"""
    <div style="background:#0a0a0a;padding:40px 20px;font-family:sans-serif;">
      <div style="max-width:500px;margin:0 auto;background:#171717;border-radius:16px;overflow:hidden;">
        <div style="background:linear-gradient(135deg,#10b981,#059669);padding:30px;text-align:center;">
          <h1 style="color:white;margin:0;">🎮 Game Invite!</h1>
        </div>
        <div style="padding:30px;">
          <p style="color:#e5e5e5;font-size:18px;"><strong style="color:#10b981;">{host_name}</strong> invited you to play!</p>
          <h2 style="color:white;text-align:center;">{group_name}</h2>
          {custom_section}
          {passcode_section}
          <a href="{invite_link}" style="display:block;background:#10b981;color:white;text-decoration:none;padding:16px;border-radius:8px;text-align:center;font-weight:bold;font-size:18px;margin-top:20px;">Join Game →</a>
          <p style="color:#525252;font-size:12px;text-align:center;margin-top:20px;">No account needed • Play as guest</p>
        </div>
      </div>
    </div>
    """
